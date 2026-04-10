"""
builders.py — Generates downloadable PowerPoint and Word files from AI-produced content.
No AI calls here — pure file generation from structured data.
"""

from io import BytesIO


# ─────────────────────────────────────────────
# COLOUR PALETTE
# ─────────────────────────────────────────────
def _rgb(r, g, b):
    from pptx.dml.color import RGBColor
    return RGBColor(r, g, b)

NAVY        = lambda: _rgb(0x1B, 0x3A, 0x6B)
WHITE       = lambda: _rgb(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = lambda: _rgb(0xF5, 0xF5, 0xF5)
MID_GRAY    = lambda: _rgb(0x6C, 0x75, 0x7D)
DARK        = lambda: _rgb(0x1A, 0x1A, 0x2E)
GREEN       = lambda: _rgb(0x28, 0xA7, 0x45)
AMBER_COLOR = lambda: _rgb(0xE6, 0xA8, 0x00)
RED_COLOR   = lambda: _rgb(0xDC, 0x35, 0x45)
AMBER_BG    = lambda: _rgb(0xFF, 0xF3, 0xCD)
AMBER_TEXT  = lambda: _rgb(0x85, 0x60, 0x04)

STATUS_COLORS = {
    "Green": lambda: _rgb(0x28, 0xA7, 0x45),
    "Amber": lambda: _rgb(0xE6, 0xA8, 0x00),
    "Red":   lambda: _rgb(0xDC, 0x35, 0x45),
}
SEVERITY_COLORS = {
    "High":   lambda: _rgb(0xDC, 0x35, 0x45),
    "Medium": lambda: _rgb(0xE6, 0xA8, 0x00),
    "Low":    lambda: _rgb(0x28, 0xA7, 0x45),
}
MS_STATUS_COLORS = {
    "Complete": lambda: _rgb(0x28, 0xA7, 0x45),
    "On Track": lambda: _rgb(0x1B, 0x3A, 0x6B),
    "At Risk":  lambda: _rgb(0xDC, 0x35, 0x45),
}


# ─────────────────────────────────────────────
# PPTX BUILDER
# ─────────────────────────────────────────────
def build_pptx(project_data: dict, slide_content: dict) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # ── Helpers ──────────────────────────────

    def set_bg(slide, color_fn):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = color_fn()

    def txt_box(slide, text, l, t, w, h,
                size=16, bold=False, italic=False,
                color_fn=DARK, align=PP_ALIGN.LEFT, wrap=True):
        from pptx.util import Inches, Pt
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = wrap
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color_fn()
        return box

    def bullet_box(slide, items, l, t, w, h,
                   size=13, color_fn=DARK, prefix="  •  "):
        from pptx.util import Inches, Pt
        box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = box.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            run = p.add_run()
            run.text = f"{prefix}{item}"
            run.font.size = Pt(size)
            run.font.color.rgb = color_fn()
            p.space_after = Pt(5)
        return box

    def header_band(slide, title, subtitle=""):
        from pptx.util import Inches, Pt
        band = slide.shapes.add_textbox(
            Inches(0), Inches(0), prs.slide_width, Inches(1.25)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = NAVY()
        band.line.fill.background()
        tf = band.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run()
        r.text = "    " + title.upper()
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = WHITE()
        if subtitle:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.LEFT
            r2 = p2.add_run()
            r2.text = "    " + subtitle
            r2.font.size = Pt(11)
            r2.font.color.rgb = _rgb(0xBB, 0xCC, 0xEE)

    def metric_card(slide, label, value, l, t):
        from pptx.util import Inches, Pt
        card = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(3.2), Inches(1.05))
        card.fill.solid()
        card.fill.fore_color.rgb = LIGHT_GRAY()
        card.line.fill.background()
        tf = card.text_frame
        tf.word_wrap = False
        p1 = tf.paragraphs[0]
        r1 = p1.add_run()
        r1.text = "  " + label.upper()
        r1.font.size = Pt(9)
        r1.font.bold = True
        r1.font.color.rgb = MID_GRAY()
        p2 = tf.add_paragraph()
        r2 = p2.add_run()
        r2.text = "  " + value
        r2.font.size = Pt(14)
        r2.font.bold = True
        r2.font.color.rgb = NAVY()

    # ── Slide data ───────────────────────────
    proj_name = slide_content.get("project_name", project_data.get("project_name", "Project Update"))
    period    = slide_content.get("reporting_period", project_data.get("period", ""))
    status    = slide_content.get("overall_status", "Green")
    client    = project_data.get("client", "")
    sc        = STATUS_COLORS.get(status, AMBER_COLOR)
    status_label = {"Green": "●  ON TRACK", "Amber": "●  AT RISK", "Red": "●  CRITICAL"}.get(status, "●  AT RISK")

    # ── SLIDE 1: Title ───────────────────────
    s1 = prs.slides.add_slide(blank)
    set_bg(s1, NAVY)

    # Accent line
    line = s1.shapes.add_textbox(Inches(1), Inches(2.6), Inches(5), Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = sc()
    line.line.fill.background()

    txt_box(s1, proj_name, 1, 2.8, 11, 1.6, size=34, bold=True, color_fn=WHITE)
    txt_box(s1, period,    1, 4.5, 8,  0.6, size=16, color_fn=lambda: _rgb(0xBB, 0xCC, 0xEE))
    txt_box(s1, status_label, 1, 5.3, 5, 0.6, size=14, bold=True, color_fn=sc)
    if client:
        txt_box(s1, client, 8, 6.5, 4.8, 0.6, size=12,
                color_fn=lambda: _rgb(0x99, 0xAA, 0xCC), align=PP_ALIGN.RIGHT)

    # ── SLIDE 2: Executive Summary ───────────
    s2 = prs.slides.add_slide(blank)
    header_band(s2, "Executive Summary", f"{period}  ·  {proj_name}")

    summary = slide_content.get("exec_summary", "")
    txt_box(s2, summary, 0.5, 1.45, 8.7, 2.8, size=14, wrap=True, color_fn=DARK)

    metrics = slide_content.get("key_metrics", [])
    my = 1.45
    for m in metrics[:4]:
        metric_card(s2, m.get("label", ""), m.get("value", "—"), 9.7, my)
        my += 1.28

    # ── SLIDE 3: Progress ────────────────────
    s3 = prs.slides.add_slide(blank)
    header_band(s3, "Progress This Period", proj_name)
    bullet_box(s3, slide_content.get("accomplishments", []), 0.5, 1.45, 12.3, 5.5,
               size=14, color_fn=DARK)

    # ── SLIDE 4: Risks & Issues ──────────────
    s4 = prs.slides.add_slide(blank)
    header_band(s4, "Risks & Issues", proj_name)

    risks = slide_content.get("risks", [])
    if risks:
        ry = 1.45
        for risk in risks[:4]:
            sev      = risk.get("severity", "Medium")
            sev_c    = SEVERITY_COLORS.get(sev, AMBER_COLOR)
            sev_text = {"High": "WHITE", "Medium": "DARK", "Low": "WHITE"}.get(sev, "WHITE")

            badge = s4.shapes.add_textbox(Inches(0.4), Inches(ry), Inches(1.1), Inches(0.45))
            badge.fill.solid()
            badge.fill.fore_color.rgb = sev_c()
            badge.line.fill.background()
            btf = badge.text_frame
            bp = btf.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = sev.upper()
            br.font.size = Pt(9)
            br.font.bold = True
            br.font.color.rgb = WHITE() if sev != "Medium" else DARK()

            txt_box(s4, risk.get("risk", ""), 1.7, ry, 11, 0.45,
                    size=13, bold=True, color_fn=DARK)
            txt_box(s4, "Mitigation: " + risk.get("mitigation", ""), 1.7, ry + 0.48,
                    11, 0.45, size=11, italic=True, color_fn=MID_GRAY)
            ry += 1.2
    else:
        txt_box(s4, "No risks or issues identified this period.", 0.5, 2.2, 12, 1,
                size=15, italic=True, color_fn=GREEN)

    # ── SLIDE 5: Milestones & Next Steps ─────
    s5 = prs.slides.add_slide(blank)
    header_band(s5, "Milestones & Next Steps", proj_name)

    txt_box(s5, "MILESTONES", 0.4, 1.45, 6, 0.35, size=10, bold=True, color_fn=MID_GRAY)
    milestones = slide_content.get("milestones", [])
    ms_y = 1.9
    for ms in milestones[:5]:
        ms_s  = ms.get("status", "On Track")
        ms_c  = MS_STATUS_COLORS.get(ms_s, lambda: _rgb(0x1B, 0x3A, 0x6B))
        icon  = {"Complete": "✓", "On Track": "→", "At Risk": "⚠"}.get(ms_s, "→")
        txt_box(s5, f"{icon}  {ms.get('name', '')}", 0.4, ms_y, 5.8, 0.42,
                size=12, bold=(ms_s == "Complete"), color_fn=ms_c)
        txt_box(s5, ms.get("date", ""), 0.4, ms_y + 0.42, 5.8, 0.32,
                size=10, italic=True, color_fn=MID_GRAY)
        ms_y += 0.9

    # Divider
    div = s5.shapes.add_textbox(Inches(6.9), Inches(1.4), Inches(0.04), Inches(5.6))
    div.fill.solid()
    div.fill.fore_color.rgb = _rgb(0xDD, 0xDD, 0xDD)
    div.line.fill.background()

    txt_box(s5, "NEXT STEPS", 7.2, 1.45, 5.7, 0.35, size=10, bold=True, color_fn=MID_GRAY)
    bullet_box(s5, slide_content.get("next_steps", []), 7.2, 1.9, 5.7, 5,
               size=13, color_fn=DARK, prefix="→  ")

    # ── SLIDE 6: Decisions Needed (conditional) ──
    decisions = slide_content.get("decisions_needed", [])
    if decisions:
        s6 = prs.slides.add_slide(blank)
        header_band(s6, "Decisions Needed", "Action Required from Leadership")
        dy = 1.55
        for dec in decisions[:4]:
            box = s6.shapes.add_textbox(Inches(0.5), Inches(dy), Inches(12.3), Inches(0.9))
            box.fill.solid()
            box.fill.fore_color.rgb = AMBER_BG()
            box.line.fill.background()
            tf = box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = f"  ⚠  {dec}"
            r.font.size = Pt(14)
            r.font.bold = True
            r.font.color.rgb = AMBER_TEXT()
            dy += 1.1

    # ── Save ─────────────────────────────────
    out = BytesIO()
    prs.save(out)
    out.seek(0)
    return out.getvalue()


# ─────────────────────────────────────────────
# DOCX BUILDER
# ─────────────────────────────────────────────
def build_docx(project_data: dict, document_content: dict, slide_content: dict) -> bytes:
    from docx import Document
    from docx.shared import Inches, Pt, RGBColor
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # Margins
    for sec in doc.sections:
        sec.top_margin    = Inches(1)
        sec.bottom_margin = Inches(1)
        sec.left_margin   = Inches(1.2)
        sec.right_margin  = Inches(1.2)

    # Title
    title = doc.add_heading(project_data.get("project_name", "Project"), 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Subtitle
    sub = doc.add_paragraph()
    sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = sub.add_run(f"Project Status Report  ·  {project_data.get('period', '')}")
    sr.font.size = Pt(12)
    sr.font.color.rgb = RGBColor(0x6C, 0x75, 0x7D)

    doc.add_paragraph()

    # Summary table
    table = doc.add_table(rows=2, cols=4)
    table.style = "Table Grid"
    headers = ["Client", "Overall Status", "Budget", "Reporting Period"]
    values  = [
        project_data.get("client", "—"),
        slide_content.get("overall_status", "—"),
        project_data.get("budget", "—"),
        project_data.get("period", "—"),
    ]
    hdr_row = table.rows[0]
    val_row = table.rows[1]
    for i, (h, v) in enumerate(zip(headers, values)):
        hdr_cell = hdr_row.cells[i]
        hdr_cell.text = h
        hdr_cell.paragraphs[0].runs[0].font.bold = True
        hdr_cell.paragraphs[0].runs[0].font.size = Pt(9)
        val_row.cells[i].text = v

    doc.add_paragraph()

    # Narrative sections
    sections = [
        ("1. Executive Summary",    document_content.get("executive_summary", "")),
        ("2. Progress This Period",  document_content.get("progress_narrative", "")),
        ("3. Risks & Issues",        document_content.get("risks_narrative", "")),
        ("4. Plan for Next Period",  document_content.get("next_period_plan", "")),
    ]
    for heading, body in sections:
        doc.add_heading(heading, level=1)
        if body:
            doc.add_paragraph(body)
        doc.add_paragraph()

    # Appendix: accomplishments
    accs = slide_content.get("accomplishments", [])
    if accs:
        doc.add_heading("Appendix: Accomplishments This Period", level=2)
        for item in accs:
            doc.add_paragraph(item, style="List Bullet")

    # Appendix: action items / next steps
    nxt = slide_content.get("next_steps", [])
    if nxt:
        doc.add_heading("Appendix: Next Steps", level=2)
        for item in nxt:
            doc.add_paragraph(item, style="List Bullet")

    out = BytesIO()
    doc.save(out)
    out.seek(0)
    return out.getvalue()
